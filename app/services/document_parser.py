"""Parse various document formats into plain text for RAG ingestion.

Supports:
- Text: .txt, .md, .csv, .tsv, .log, .json, .xml, .yaml, .yml, .rst
- PDF: .pdf
- Word: .docx, .doc (via python-docx)
- Excel: .xlsx, .xls (via openpyxl)
- PowerPoint: .pptx, .ppt (via python-pptx)
- OpenDocument: .odt, .ods, .odp (via odfpy)
- Rich Text: .rtf

Returns plain text suitable for chunking and embedding.
"""

from __future__ import annotations

import csv
import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# Extensions we can read as plain text directly
TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".tsv", ".log", ".json", ".xml",
    ".yaml", ".yml", ".rst", ".text", ".markdown", ".rtf",
    ".ini", ".cfg", ".conf", ".env", ".sh", ".py", ".js",
    ".html", ".htm", ".css",
}


def extract_text(file_bytes: bytes, filename: str) -> str:
    """Extract plain text from a file based on its extension.

    Raises ValueError if the format is not supported.
    """
    ext = Path(filename).suffix.lower()

    # Plain text formats
    if ext in TEXT_EXTENSIONS:
        return file_bytes.decode("utf-8", errors="replace")

    # PDF
    if ext == ".pdf":
        return _parse_pdf(file_bytes)

    # Word
    if ext in (".docx", ".doc"):
        return _parse_docx(file_bytes)

    # Excel
    if ext in (".xlsx", ".xls"):
        return _parse_xlsx(file_bytes)

    # PowerPoint
    if ext in (".pptx", ".ppt"):
        return _parse_pptx(file_bytes)

    # OpenDocument Text
    if ext == ".odt":
        return _parse_odt(file_bytes)

    # OpenDocument Spreadsheet
    if ext == ".ods":
        return _parse_ods(file_bytes)

    # OpenDocument Presentation
    if ext == ".odp":
        return _parse_odp(file_bytes)

    raise ValueError(
        f"Формат «{ext}» не поддерживается. "
        f"Поддерживаемые форматы: txt, md, csv, pdf, docx, xlsx, pptx, odt, ods, odp"
    )


def _parse_pdf(data: bytes) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(data))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)
    except ImportError:
        raise ValueError("PDF parsing requires PyPDF2. Install: pip install PyPDF2")
    except Exception as e:
        raise ValueError(f"Не удалось прочитать PDF: {e}")


def _parse_docx(data: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n\n".join(paragraphs)
    except ImportError:
        raise ValueError("DOCX parsing requires python-docx. Install: pip install python-docx")
    except Exception as e:
        raise ValueError(f"Не удалось прочитать DOCX: {e}")


def _parse_xlsx(data: bytes) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), data_only=True)
        lines = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            lines.append(f"=== {sheet} ===")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c for c in cells):
                    lines.append(" | ".join(cells))
        return "\n".join(lines)
    except ImportError:
        raise ValueError("XLSX parsing requires openpyxl. Install: pip install openpyxl")
    except Exception as e:
        raise ValueError(f"Не удалось прочитать XLSX: {e}")


def _parse_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides.append(f"--- Слайд {i} ---\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except ImportError:
        raise ValueError("PPTX parsing requires python-pptx. Install: pip install python-pptx")
    except Exception as e:
        raise ValueError(f"Не удалось прочитать PPTX: {e}")


def _parse_odt(data: bytes) -> str:
    try:
        from odf.opendocument import load
        from odf.text import P
        doc = load(io.BytesIO(data))
        paragraphs = []
        for p in doc.getElementsByType(P):
            text = ""
            for node in p.childNodes:
                if hasattr(node, "data"):
                    text += node.data
            if text.strip():
                paragraphs.append(text.strip())
        return "\n\n".join(paragraphs)
    except ImportError:
        raise ValueError("ODT parsing requires odfpy. Install: pip install odfpy")
    except Exception as e:
        raise ValueError(f"Не удалось прочитать ODT: {e}")


def _parse_ods(data: bytes) -> str:
    try:
        from odf.opendocument import load
        from odf.table import Table, TableRow, TableCell
        from odf.text import P
        doc = load(io.BytesIO(data))
        lines = []
        for table in doc.getElementsByType(Table):
            name = table.getAttribute("name") or "Sheet"
            lines.append(f"=== {name} ===")
            for row in table.getElementsByType(TableRow):
                cells = []
                for cell in row.getElementsByType(TableCell):
                    text = ""
                    for p in cell.getElementsByType(P):
                        for node in p.childNodes:
                            if hasattr(node, "data"):
                                text += node.data
                    cells.append(text.strip())
                if any(c for c in cells):
                    lines.append(" | ".join(cells))
        return "\n".join(lines)
    except Exception as e:
        raise ValueError(f"Не удалось прочитать ODS: {e}")


def _parse_odp(data: bytes) -> str:
    try:
        from odf.opendocument import load
        from odf.draw import Frame, Page
        from odf.text import P
        doc = load(io.BytesIO(data))
        slides = []
        for i, page in enumerate(doc.getElementsByType(Page), 1):
            texts = []
            for p in page.getElementsByType(P):
                text = ""
                for node in p.childNodes:
                    if hasattr(node, "data"):
                        text += node.data
                if text.strip():
                    texts.append(text.strip())
            if texts:
                slides.append(f"--- Слайд {i} ---\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        raise ValueError(f"Не удалось прочитать ODP: {e}")
